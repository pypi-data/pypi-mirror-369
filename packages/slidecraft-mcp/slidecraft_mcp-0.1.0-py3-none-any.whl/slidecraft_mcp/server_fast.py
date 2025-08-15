import asyncio
import tempfile
import os
import uuid
from typing import Dict, Any, Optional, List
import base64
from pathlib import Path

from mcp.server.models import InitializationOptions
import mcp.types as types
from mcp.server import NotificationOptions, Server
from pydantic import AnyUrl
import mcp.server.stdio

# Store active presentations - minimal storage for fast startup
presentations: Dict[str, Dict[str, Any]] = {}
presentation_metadata: Dict[str, Dict[str, Any]] = {}

server = Server("powerpoint-mcp-server")

@server.list_tools()
async def handle_list_tools() -> list[types.Tool]:
    """
    List available PowerPoint tools.
    Each tool specifies its arguments using JSON Schema.
    """
    return [
        types.Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Presentation title"
                    },
                    "template": {
                        "type": "string",
                        "description": "Template to use (blank, title_slide, etc.)",
                        "default": "blank"
                    },
                    "slide_size": {
                        "type": "string",
                        "enum": ["widescreen", "standard"],
                        "description": "Slide dimensions",
                        "default": "widescreen"
                    }
                },
                "required": ["title"]
            }
        ),
        types.Tool(
            name="add_slide",
            description="Add a new slide to presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "Presentation ID"
                    },
                    "layout": {
                        "type": "string",
                        "description": "Slide layout type",
                        "default": "blank"
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title"
                    }
                },
                "required": ["presentation_id"]
            }
        ),
        types.Tool(
            name="add_text",
            description="Add text to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "Presentation ID"
                    },
                    "slide_index": {
                        "type": "integer",
                        "description": "Slide index (0-based)"
                    },
                    "text": {
                        "type": "string",
                        "description": "Text content"
                    },
                    "position": {
                        "type": "object",
                        "description": "Position and size {x, y, width, height} in inches"
                    }
                },
                "required": ["presentation_id", "slide_index", "text"]
            }
        ),
        types.Tool(
            name="save_presentation",
            description="Save presentation to file",
            inputSchema={
                "type": "object",
                "properties": {
                    "presentation_id": {
                        "type": "string",
                        "description": "Presentation ID"
                    },
                    "filename": {
                        "type": "string",
                        "description": "Output filename"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["pptx", "pdf"],
                        "description": "Export format",
                        "default": "pptx"
                    }
                },
                "required": ["presentation_id", "filename"]
            }
        ),
        types.Tool(
            name="list_presentations",
            description="List all active presentations",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        )
    ]

@server.call_tool()
async def handle_call_tool(name: str, arguments: dict[str, Any]) -> list[types.TextContent]:
    """Handle tool calls by routing to appropriate handlers."""
    try:
        if name == "create_presentation":
            return await handle_create_presentation(arguments)
        elif name == "add_slide":
            return await handle_add_slide(arguments)
        elif name == "add_text":
            return await handle_add_text(arguments)
        elif name == "save_presentation":
            return await handle_save_presentation(arguments)
        elif name == "list_presentations":
            return await handle_list_presentations(arguments)
        else:
            return [types.TextContent(type="text", text=f"Unknown tool: {name}")]
    except Exception as e:
        return [types.TextContent(type="text", text=f"Error: {str(e)}")]

async def handle_create_presentation(arguments: dict[str, Any]) -> list[types.TextContent]:
    """Create a new PowerPoint presentation with lazy loading."""
    try:
        # Lazy import - only load when actually needed
        from pptx import Presentation
        from pptx.util import Inches
        
        title = arguments.get("title", "New Presentation")
        template = arguments.get("template", "blank")
        slide_size = arguments.get("slide_size", "widescreen")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size
        if slide_size == "widescreen":
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
        else:  # standard
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # Generate unique ID
        presentation_id = str(uuid.uuid4())
        
        # Store presentation (store the actual object for now)
        presentations[presentation_id] = {"prs": prs}
        presentation_metadata[presentation_id] = {
            "title": title,
            "template": template,
            "slide_size": slide_size,
            "created_at": "now",
            "slide_count": len(prs.slides)
        }
        
        return [types.TextContent(
            type="text",
            text=f"✅ Created presentation '{title}' with ID: {presentation_id}"
        )]
    except ImportError as e:
        return [types.TextContent(
            type="text",
            text=f"❌ Missing dependency: {str(e)}. Please install python-pptx"
        )]
    except Exception as e:
        return [types.TextContent(
            type="text",
            text=f"❌ Error creating presentation: {str(e)}"
        )]

async def handle_add_slide(arguments: dict[str, Any]) -> list[types.TextContent]:
    """Add a new slide to presentation."""
    try:
        presentation_id = arguments.get("presentation_id")
        layout = arguments.get("layout", "blank")
        title = arguments.get("title", "")
        
        if presentation_id not in presentations:
            return [types.TextContent(
                type="text",
                text=f"❌ Presentation {presentation_id} not found"
            )]
        
        prs = presentations[presentation_id]["prs"]
        
        # Add slide with blank layout (index 6 is typically blank)
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Update metadata
        presentation_metadata[presentation_id]["slide_count"] = len(prs.slides)
        
        slide_index = len(prs.slides) - 1
        
        return [types.TextContent(
            type="text",
            text=f"✅ Added slide {slide_index} to presentation {presentation_id}"
        )]
    except Exception as e:
        return [types.TextContent(
            type="text",
            text=f"❌ Error adding slide: {str(e)}"
        )]

async def handle_add_text(arguments: dict[str, Any]) -> list[types.TextContent]:
    """Add text to a slide."""
    try:
        from pptx.util import Inches, Pt
        
        presentation_id = arguments.get("presentation_id")
        slide_index = arguments.get("slide_index", 0)
        text = arguments.get("text", "")
        position = arguments.get("position", {"x": 1, "y": 1, "width": 8, "height": 1})
        
        if presentation_id not in presentations:
            return [types.TextContent(
                type="text",
                text=f"❌ Presentation {presentation_id} not found"
            )]
        
        prs = presentations[presentation_id]["prs"]
        
        if slide_index >= len(prs.slides):
            return [types.TextContent(
                type="text",
                text=f"❌ Slide index {slide_index} out of range"
            )]
        
        slide = prs.slides[slide_index]
        
        # Add text box
        left = Inches(position.get("x", 1))
        top = Inches(position.get("y", 1))
        width = Inches(position.get("width", 8))
        height = Inches(position.get("height", 1))
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        
        return [types.TextContent(
            type="text",
            text=f"✅ Added text to slide {slide_index}: '{text[:50]}...'"
        )]
    except Exception as e:
        return [types.TextContent(
            type="text",
            text=f"❌ Error adding text: {str(e)}"
        )]

async def handle_save_presentation(arguments: dict[str, Any]) -> list[types.TextContent]:
    """Save presentation to file."""
    try:
        presentation_id = arguments.get("presentation_id")
        filename = arguments.get("filename")
        format_type = arguments.get("format", "pptx")
        
        if presentation_id not in presentations:
            return [types.TextContent(
                type="text",
                text=f"❌ Presentation {presentation_id} not found"
            )]
        
        prs = presentations[presentation_id]["prs"]
        
        # Ensure filename has correct extension
        if not filename.endswith(f".{format_type}"):
            filename = f"{filename}.{format_type}"
        
        # Save the file
        if format_type == "pptx":
            prs.save(filename)
            file_size = os.path.getsize(filename)
            
            return [types.TextContent(
                type="text",
                text=f"✅ Saved presentation to {filename} ({file_size} bytes)"
            )]
        else:
            return [types.TextContent(
                type="text",
                text=f"❌ Format {format_type} not yet supported. Only PPTX is available."
            )]
    except Exception as e:
        return [types.TextContent(
            type="text",
            text=f"❌ Error saving presentation: {str(e)}"
        )]

async def handle_list_presentations(arguments: dict[str, Any]) -> list[types.TextContent]:
    """List all active presentations."""
    if not presentations:
        return [types.TextContent(
            type="text",
            text="📝 No active presentations"
        )]
    
    result = "📋 Active Presentations:\n\n"
    for pid, metadata in presentation_metadata.items():
        result += f"🆔 ID: {pid}\n"
        result += f"📄 Title: {metadata['title']}\n"
        result += f"📊 Slides: {metadata['slide_count']}\n"
        result += f"📏 Size: {metadata['slide_size']}\n"
        result += f"⏰ Created: {metadata['created_at']}\n\n"
    
    return [types.TextContent(type="text", text=result)]

async def main():
    # Run the server using stdin/stdout streams
    async with mcp.server.stdio.stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="powerpoint-mcp-server",
                server_version="0.1.0",
                capabilities=server.get_capabilities(
                    notification_options=NotificationOptions(),
                    experimental_capabilities={},
                ),
            ),
        )

if __name__ == "__main__":
    asyncio.run(main())
