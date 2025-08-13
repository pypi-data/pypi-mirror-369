from . import register_extractor
import io
import zipfile
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

if Presentation:
 def xtxt_pptx(file_buffer) -> str:
    try:
        # Convertiamo il file_buffer (che è già un BytesIO o simile) in modo da poterlo riusare
        file_buffer.seek(0)
        data = file_buffer.read()
        buffer_copy = io.BytesIO(data)

        if not zipfile.is_zipfile(buffer_copy):
            print("⚠️  Invalid PPTX (not a ZIP file)" )
            return ""

        # Se è un file zip valido, possiamo ripassare i dati a Presentation
        buffer_copy.seek(0)
        prs = Presentation(buffer_copy)

        text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        return text

    except Exception as e:
        print(f"⚠️ Error during PPTX extraction: {e}")
        return ""
 register_extractor(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation" ,
    xtxt_pptx,
    name="PPTX"
)
